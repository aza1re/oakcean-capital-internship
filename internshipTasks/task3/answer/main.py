# ...existing code...
import time
import csv
from statistics import mean
import requests
import os
from pptx import Presentation
from pptx.util import Inches

API_URL = "http://127.0.0.1:8000/timeseries"  # adjust if needed
OUTPUT_CSV = "task3_benchmark.csv"
PPTX_OUT = "task3_benchmark_summary.pptx"

QUERIES = [
    {"tickers": "RELIANCE,TCS", "start": "2023-01-02T09:15:00", "end": "2023-01-02T15:30:00", "fields": "open,close,volume"},
    {"tickers": "RELIANCE,TCS,INFY", "start": "2023-02-01T09:15:00", "end": "2023-02-01T15:30:00", "fields": "open,high,low,close,volume"},
]

DBS = ["mongo", "influx"]

def run_benchmark(repeats: int = 5, timeout: int = 60):
    rows = []
    for db in DBS:
        for q in QUERIES:
            params = q.copy()
            params["db"] = db
            latencies = []
            success_counts = 0
            for i in range(repeats):
                try:
                    t0 = time.time()
                    r = requests.get(API_URL, params=params, timeout=timeout)
                    elapsed = time.time() - t0
                    latencies.append(elapsed)
                    if r.status_code == 200:
                        success_counts += 1
                except Exception as e:
                    latencies.append(None)
                time.sleep(0.2)
            valid = [l for l in latencies if l is not None]
            mean_latency = mean(valid) if valid else None
            rows.append({
                "db": db,
                "tickers": params["tickers"],
                "start": params["start"],
                "end": params["end"],
                "fields": params["fields"],
                "latencies": latencies,
                "mean": mean_latency,
                "success_count": success_counts
            })
    # write CSV
    with open(OUTPUT_CSV, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["db", "tickers", "start", "end", "fields", "latencies", "mean", "success_count"])
        for r in rows:
            writer.writerow([r["db"], r["tickers"], r["start"], r["end"], r["fields"], ";".join([str(x) for x in r["latencies"]]), r["mean"], r["success_count"]])
    # create simple pptx
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]
    for r in rows:
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f'Benchmark: {r["db"]} — {r["tickers"]}'
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
        tf = body.text_frame
        tf.text = f'Mean latency: {r["mean"]:.4f} s' if r["mean"] else "Mean latency: N/A"
        p = tf.add_paragraph()
        p.text = f'Runs: {r["latencies"]}'
        p.level = 1
        p2 = tf.add_paragraph()
        p2.text = f'Successful responses: {r["success_count"]}'
        p2.level = 1
    prs.save(PPTX_OUT)
    print(f"Saved CSV: {os.path.abspath(OUTPUT_CSV)}, PPTX: {os.path.abspath(PPTX_OUT)}")

if __name__ == "__main__":
    run_benchmark()
# ...existing code...